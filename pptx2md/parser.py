# Copyright 2024 Liu Siyao
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#
# Modifications Copyright 2025-2026 vanilla1108

from __future__ import print_function

import logging
import os
import shutil
import subprocess
import sys
from functools import partial
from operator import attrgetter
from typing import List, Tuple, Union

from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from rapidfuzz import process as fuze_process
from tqdm import tqdm

from pptx2md.types import (
    ConversionConfig,
    GeneralSlide,
    ImageElement,
    ListItemElement,
    ListType,
    ParagraphElement,
    ParsedPresentation,
    SlideElement,
    TableElement,
    TextRun,
    TextStyle,
    TitleElement,
)

logger = logging.getLogger(__name__)

picture_count = 0
_SLIDE_SIZE_EMU_CACHE: dict[str, tuple[int, int]] = {}
_PPT_COM_SESSION = None
_WMF_COM_FALLBACK_STATE: bool | None = None  # None=未探测/未尝试, True=可用, False=不可用/禁用
_WMF_COM_FALLBACK_REASON_LOGGED = False


class _PowerPointComSession:
    def __init__(self):
        self._app = None
        self._pres = None
        self._pptx_path = None

    def ensure_open(self, pptx_path: str):
        normalized_pptx_path = os.path.abspath(str(pptx_path))
        if self._app is not None and self._pres is not None and self._pptx_path == normalized_pptx_path:
            return
        self.close()
        import win32com.client  # type: ignore

        self._app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            # 尽量不弹窗；某些环境下 WithWindow=False 会异常，所以用 WithWindow=True + Visible=False
            self._app.Visible = False
        except Exception:
            pass
        # ReadOnly=True, Untitled=False, WithWindow=True
        self._pres = self._app.Presentations.Open(normalized_pptx_path, True, False, True)
        self._pptx_path = normalized_pptx_path

    def export_slide_png(self, slide_idx: int, slide_png_path: str, width_px: int, height_px: int) -> bool:
        if self._pres is None:
            return False
        normalized_slide_png_path = os.path.abspath(str(slide_png_path))
        try:
            com_slide = self._pres.Slides(slide_idx)  # 1-based
            com_slide.Export(normalized_slide_png_path, "PNG", width_px, height_px)
            return os.path.exists(normalized_slide_png_path) and os.path.getsize(normalized_slide_png_path) > 0
        except Exception:
            return False

    def close(self):
        try:
            if self._pres is not None:
                self._pres.Close()
        except Exception:
            pass
        try:
            if self._app is not None:
                self._app.Quit()
        except Exception:
            pass
        self._app = None
        self._pres = None
        self._pptx_path = None


def close_powerpoint_com_session():
    """关闭全局 COM 会话，避免 PowerPoint 残留后台进程。"""
    global _PPT_COM_SESSION
    try:
        if _PPT_COM_SESSION is not None:
            _PPT_COM_SESSION.close()
    finally:
        _PPT_COM_SESSION = None


def _should_try_wmf_com_fallback() -> bool:
    """是否尝试使用 PowerPoint COM 将 WMF 转为高清光栅图。

    规则：
    - 若显式设置了环境变量 PPTX2MD_WMF_COM_FALLBACK，则以其为准（0/1）。
    - 否则自动探测：在 Windows 上且能 import win32com 时，默认尝试一次；失败后缓存为 False。
    """
    global _WMF_COM_FALLBACK_STATE
    global _WMF_COM_FALLBACK_REASON_LOGGED

    def _log_disable_once(message: str):
        global _WMF_COM_FALLBACK_REASON_LOGGED
        if not _WMF_COM_FALLBACK_REASON_LOGGED:
            logger.info(message)
            _WMF_COM_FALLBACK_REASON_LOGGED = True

    # 用户显式开关优先
    env = os.environ.get("PPTX2MD_WMF_COM_FALLBACK")
    if env is not None:
        enabled = _env_bool("PPTX2MD_WMF_COM_FALLBACK", default=False)
        if not enabled:
            _log_disable_once("WMF 的 PowerPoint COM 兜底已禁用（PPTX2MD_WMF_COM_FALLBACK=0）。")
            return False
        if os.name != "nt":
            _log_disable_once("WMF 的 PowerPoint COM 兜底仅在 Windows 可用，当前环境已跳过。")
            return False
        try:
            import win32com.client  # noqa: F401

            return True
        except Exception:
            _log_disable_once("未检测到 win32com，无法启用 WMF 的 PowerPoint COM 兜底。")
            return False

    if _WMF_COM_FALLBACK_STATE is not None:
        return _WMF_COM_FALLBACK_STATE

    # pytest 环境默认禁用 COM，避免测试进程因 PowerPoint/COM 生命周期不稳定而崩溃
    if os.environ.get("PYTEST_CURRENT_TEST") is not None or "pytest" in sys.modules:
        _WMF_COM_FALLBACK_STATE = False
        _log_disable_once("检测到 pytest 运行环境，WMF 的 PowerPoint COM 兜底默认禁用。")
        return False

    if os.name != "nt":
        _WMF_COM_FALLBACK_STATE = False
        _log_disable_once("当前不是 Windows，WMF 的 PowerPoint COM 兜底不可用。")
        return False

    try:
        import win32com.client  # noqa: F401

        _WMF_COM_FALLBACK_STATE = True
        return True
    except Exception:
        _WMF_COM_FALLBACK_STATE = False
        _log_disable_once("未检测到 win32com，WMF 的 PowerPoint COM 兜底不可用。")
        return False


def _get_slide_size_emu(pptx_path) -> tuple[int, int]:
    """读取演示文稿的页面尺寸（EMU），用于将 shape 边界映射到导出的像素坐标。"""
    key = os.path.abspath(str(pptx_path))
    cached = _SLIDE_SIZE_EMU_CACHE.get(key)
    if cached:
        return cached
    prs = Presentation(pptx_path)
    size = (int(prs.slide_width), int(prs.slide_height))
    _SLIDE_SIZE_EMU_CACHE[key] = size
    return size


def _env_int(name: str, default: int) -> int:
    try:
        v = int(os.environ.get(name, str(default)).strip())
        return v
    except Exception:
        return default


def _env_bool(name: str, default: bool = False) -> bool:
    v = os.environ.get(name)
    if v is None:
        return default
    return v.strip().lower() in ("1", "true", "yes", "y", "on")


def _env_str(name: str, default: str) -> str:
    v = os.environ.get(name)
    if v is None:
        return default
    return v.strip()


def _convert_wmf_via_magick(input_wmf_path: str, output_path: str, dpi: int, jpg_quality: int) -> bool:
    """用 ImageMagick CLI 将 WMF 光栅化为 PNG/JPG（需要本机安装 magick 且支持 WMF）。"""
    magick = shutil.which("magick")
    if not magick:
        return False
    _, out_ext = os.path.splitext(output_path)
    out_ext = out_ext.lower().lstrip(".")
    cmd = [
        magick,
        "-density",
        str(dpi),
        input_wmf_path,
        "-background",
        "white",
        "-alpha",
        "remove",
        "-alpha",
        "off",
    ]
    if out_ext in ("jpg", "jpeg"):
        cmd += ["-quality", str(jpg_quality)]
    cmd.append(output_path)
    try:
        subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return os.path.exists(output_path) and os.path.getsize(output_path) > 0
    except Exception:
        return False


def _convert_wmf_via_powerpoint_slide_export(config: ConversionConfig, shape, slide_idx: int, output_path) -> bool:
    """用 PowerPoint COM 将整页导出为 PNG，再按 shape 边界裁剪得到高清图片。

    这是针对 WMF 解码器缺失时的兜底方案（不依赖 ImageMagick/libwmf）。
    """
    try:
        global _PPT_COM_SESSION
        global _WMF_COM_FALLBACK_STATE
        slide_w_emu, slide_h_emu = _get_slide_size_emu(config.pptx_path)
        export_width_px = _env_int("PPTX2MD_WMF_COM_EXPORT_WIDTH", 3840)
        export_height_px = max(1, int(round(export_width_px * slide_h_emu / slide_w_emu)))

        # 复用整页导出，避免同一页多个 WMF 重复启动 PowerPoint
        slide_export_dir = config.image_dir / "_slide_exports"
        os.makedirs(slide_export_dir, exist_ok=True)
        slide_png = slide_export_dir / f"slide_{slide_idx}_w{export_width_px}.png"

        if not slide_png.exists():
            try:
                if _PPT_COM_SESSION is None:
                    _PPT_COM_SESSION = _PowerPointComSession()
                _PPT_COM_SESSION.ensure_open(os.path.abspath(str(config.pptx_path)))
                _PPT_COM_SESSION.export_slide_png(
                    slide_idx,
                    os.path.abspath(str(slide_png)),
                    export_width_px,
                    export_height_px,
                )
            except Exception as e:
                logger.warning(f'PowerPoint COM fallback failed in slide {slide_idx}: {e}')
                # 自动探测模式下，如果 COM 初始化失败，后续不再尝试
                if os.environ.get("PPTX2MD_WMF_COM_FALLBACK") is None:
                    _WMF_COM_FALLBACK_STATE = False
                return False

        if not slide_png.exists():
            return False

        # EMU -> pixel 映射后裁剪
        scale = export_width_px / float(slide_w_emu)
        left = int(round(float(shape.left) * scale))
        top = int(round(float(shape.top) * scale))
        right = int(round((float(shape.left) + float(shape.width)) * scale))
        bottom = int(round((float(shape.top) + float(shape.height)) * scale))

        # 裁剪安全边界
        with Image.open(slide_png) as slide_img:
            w, h = slide_img.size
            left = max(0, min(left, w - 1))
            top = max(0, min(top, h - 1))
            right = max(left + 1, min(right, w))
            bottom = max(top + 1, min(bottom, h))

            cropped = slide_img.crop((left, top, right, bottom))
            ext = str(os.path.splitext(str(output_path))[1]).lower()
            if ext in (".jpg", ".jpeg"):
                cropped = cropped.convert("RGB")
                cropped.save(output_path, quality=_env_int("PPTX2MD_WMF_JPEG_QUALITY", 92), optimize=True)
            else:
                cropped.save(output_path)

        return os.path.exists(output_path) and os.path.getsize(output_path) > 0
    except Exception:
        if os.environ.get("PPTX2MD_WMF_COM_FALLBACK") is None:
            _WMF_COM_FALLBACK_STATE = False
        return False


def is_title(shape):
    if shape.is_placeholder and (shape.placeholder_format.type == PP_PLACEHOLDER.TITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.VERTICAL_TITLE or
                                 shape.placeholder_format.type == PP_PLACEHOLDER.CENTER_TITLE):
        return True
    return False


_NS_A14 = '{http://schemas.microsoft.com/office/drawing/2010/main}'


def is_text_block(config: ConversionConfig, shape):
    if shape.has_text_frame:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            return True
        if len(shape.text) > config.min_block_size:
            return True
        # 含公式的 shape（a14:m）也应被当作文本块处理，
        # 因为 shape.text 不包含公式内容，会低估文本长度
        if any(True for _ in shape._element.iter(f'{_NS_A14}m')):
            return True
    return False


def _check_bullet_in_ppr(pPr):
    """检查段落属性元素中的项目符号类型。

    返回: 'bullet'、'numbered'、'none' 或 None（未找到显式设置）。
    """
    if pPr is None:
        return None
    if pPr.find(qn('a:buNone')) is not None:
        return 'none'
    if pPr.find(qn('a:buChar')) is not None:
        return 'bullet'
    bu_auto = pPr.find(qn('a:buAutoNum'))
    if bu_auto is not None:
        return 'numbered'
    return None


def _check_list_style_for_level(lst_style_elem, level):
    """在 <a:lstStyle> 元素中查找指定级别的项目符号定义。"""
    if lst_style_elem is None:
        return None
    level_tag = qn(f'a:lvl{level + 1}pPr')
    lvl_ppr = lst_style_elem.find(level_tag)
    return _check_bullet_in_ppr(lvl_ppr)


def get_paragraph_bullet_type(para, shape=None):
    """检测单个段落的项目符号类型。

    检测层级：
    1. 段落自身 XML 的显式设置（buChar/buAutoNum/buNone）
    2. 文本体的 lstStyle 默认值
    3. 占位符从布局/母版继承的列表样式
    4. 备注占位符从 notes master 继承的 notesStyle（shape 可用时）
    5. 启发式：BODY 占位符中同一 shape 存在其他显式列表项时视为列表（shape 可用时）
    6. para.level > 0 → 'bullet'
    7. 默认 → 'none'

    返回: 'bullet'、'numbered' 或 'none'。
    """
    level = para.level

    # 层级1：段落自身的显式设置
    pPr = para._p.find(qn('a:pPr'))
    explicit = _check_bullet_in_ppr(pPr)
    if explicit is not None:
        return explicit

    # 层级2：文本体的列表样式
    tx_body = None
    if shape is not None:
        try:
            tx_body = shape.text_frame._txBody
        except Exception:
            tx_body = None
    lst_style = tx_body.find(qn('a:lstStyle')) if tx_body is not None else None
    body_result = _check_list_style_for_level(lst_style, level)
    if body_result is not None:
        return body_result

    # 层级3：占位符继承 —— 从布局和母版中查找列表样式
    if shape is not None and shape.is_placeholder:
        try:
            ph_idx = shape.placeholder_format.idx
            layout = shape.part.slide_layout
            for layout_ph in layout.placeholders:
                if layout_ph.placeholder_format.idx == ph_idx:
                    layout_lst = layout_ph.text_frame._txBody.find(qn('a:lstStyle'))
                    layout_result = _check_list_style_for_level(layout_lst, level)
                    if layout_result is not None:
                        return layout_result
                    break
        except Exception:
            pass
        try:
            ph_idx = shape.placeholder_format.idx
            master = shape.part.slide_layout.slide_master
            for master_ph in master.placeholders:
                if master_ph.placeholder_format.idx == ph_idx:
                    master_lst = master_ph.text_frame._txBody.find(qn('a:lstStyle'))
                    master_result = _check_list_style_for_level(master_lst, level)
                    if master_result is not None:
                        return master_result
                    break
        except Exception:
            pass

    # 层级4：备注占位符继承 —— 从 notes master 中查找 notesStyle
    if shape is not None and shape.is_placeholder:
        try:
            notes_master = shape.part.notes_master
            notes_style = notes_master.element.find(qn('p:notesStyle'))
            notes_result = _check_list_style_for_level(notes_style, level)
            if notes_result is not None:
                return notes_result
        except Exception:
            pass

    # 层级5：level > 0 通常意味着列表项
    if level > 0:
        return 'bullet'

    return 'none'


def get_paragraph_number_start_at(para) -> int | None:
    """读取段落 buAutoNum 的显式 startAt（若未显式设置则返回 None）。"""
    pPr = para._p.find(qn('a:pPr'))
    if pPr is None:
        return None

    bu_auto = pPr.find(qn('a:buAutoNum'))
    if bu_auto is None:
        return None

    start_at = bu_auto.get('startAt')
    if start_at is None:
        return None

    try:
        value = int(start_at)
        return value if value > 0 else None
    except (TypeError, ValueError):
        return None


def is_accent(font):
    if font.underline or font.italic or (
            font.color.type == MSO_COLOR_TYPE.SCHEME and
        (font.color.theme_color == MSO_THEME_COLOR.ACCENT_1 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_2 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_3 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_4 or
         font.color.theme_color == MSO_THEME_COLOR.ACCENT_5 or font.color.theme_color == MSO_THEME_COLOR.ACCENT_6)):
        return True
    return False


def is_strong(font):
    if font.bold or (font.color.type == MSO_COLOR_TYPE.SCHEME and (font.color.theme_color == MSO_THEME_COLOR.DARK_1 or
                                                                   font.color.theme_color == MSO_THEME_COLOR.DARK_2)):
        return True
    return False


_NS_MATH = '{http://schemas.openxmlformats.org/officeDocument/2006/math}'


def _omml_to_latex(a14m_element) -> Tuple[str, bool]:
    """将 <a14:m> 元素中的 OMML 转为 LaTeX 字符串。

    支持两种包裹形态：
    - a14:m > m:oMath（行内公式）
    - a14:m > m:oMathPara > m:oMath（块级公式段落）

    返回 (text, is_math)。is_math=True 时 text 为 LaTeX，否则为纯文本回退。
    """
    from dwml.omml import oMath2Latex

    omath = a14m_element.find(f'{_NS_MATH}oMath')
    if omath is None:
        omathpara = a14m_element.find(f'{_NS_MATH}oMathPara')
        if omathpara is not None:
            omath = omathpara.find(f'{_NS_MATH}oMath')
    if omath is None:
        return ('[公式]', False)
    try:
        latex = str(oMath2Latex(omath)).strip()
        if latex:
            return (latex, True)
    except Exception as e:
        logger.warning(f'OMML→LaTeX 转换失败: {e}')
    # 回退：收集所有 m:t 文本
    texts = [t.text for t in omath.iter(f'{_NS_MATH}t') if t.text]
    if texts:
        return (''.join(texts), False)
    return ('[公式]', False)


def get_text_runs(para) -> List[TextRun]:
    runs = []
    r_index = 0
    pptx_runs = para.runs
    for child in para._p:
        local_tag = child.tag.split('}')[-1]
        if local_tag == 'r':
            if r_index < len(pptx_runs):
                run = pptx_runs[r_index]
                r_index += 1
                result = TextRun(text=run.text, style=TextStyle())
                if result.text == '':
                    continue
                try:
                    if run.hyperlink.address:
                        result.style.hyperlink = run.hyperlink.address
                except:
                    result.style.hyperlink = 'error:ppt-link-parsing-issue'
                if is_accent(run.font):
                    result.style.is_accent = True
                if is_strong(run.font):
                    result.style.is_strong = True
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    result.style.color_rgb = run.font.color.rgb
                runs.append(result)
        elif local_tag == 'm':
            text, is_math = _omml_to_latex(child)
            if text:
                runs.append(TextRun(text=text, style=TextStyle(is_math=is_math)))
    return runs


def process_title(config: ConversionConfig, shape, slide_idx) -> TitleElement:
    text = shape.text_frame.text.strip()
    if config.custom_titles:
        res = fuze_process.extractOne(text, config.custom_titles.keys(), score_cutoff=92)
        if not res:
            return TitleElement(content=text.strip(), level=max(config.custom_titles.values()) + 1)
        else:
            logger.info(f'Title in slide {slide_idx} "{text}" is converted to "{res[0]}" as specified in title file.')
            return TitleElement(content=res[0].strip(), level=config.custom_titles[res[0]])
    else:
        return TitleElement(content=text.strip(), level=1)


def _indent_list_items_after_paragraph(elements):
    """同一文本块内，段落后紧跟的 level 0 列表项提升一级缩进。

    仅当列表组的最小 level 为 0 时才推升——已经缩进的列表项
    （min level > 0）本身就是段落的子项，不需要额外加。
    """
    if len(elements) < 2:
        return elements
    result = list(elements)
    i = 0
    while i < len(result):
        if isinstance(result[i], ParagraphElement):
            # 找到段落后连续的列表项区间 [i+1, j)
            j = i + 1
            while j < len(result) and isinstance(result[j], ListItemElement):
                j += 1
            list_slice = result[i + 1:j]
            if list_slice:
                min_level = min(item.level for item in list_slice)
                if min_level == 0:
                    for k in range(i + 1, j):
                        result[k] = result[k].model_copy(update={'level': result[k].level + 1})
            i = j
        else:
            i += 1
    return result


def process_text_blocks(config: ConversionConfig, shape, slide_idx) -> List[Union[ListItemElement, ParagraphElement]]:
    results = []
    for para in shape.text_frame.paragraphs:
        has_text = para.text.strip() != ''
        has_math = any(c.tag.split('}')[-1] == 'm' for c in para._p)
        if not has_text and not has_math:
            continue
        text = get_text_runs(para)
        bullet_type = get_paragraph_bullet_type(para, shape)
        if bullet_type == 'numbered':
            results.append(
                ListItemElement(
                    content=text,
                    level=para.level,
                    list_type=ListType.Ordered,
                    list_number=get_paragraph_number_start_at(para),
                )
            )
        elif bullet_type == 'bullet':
            results.append(ListItemElement(content=text, level=para.level, list_type=ListType.Unordered))
        else:
            results.append(ParagraphElement(content=text))
    return _indent_list_items_after_paragraph(results)


def _get_notes_text_shape(notes_slide):
    """定位备注页文本框所属 shape，用于复用列表样式判定。"""
    try:
        note_shape = notes_slide.notes_placeholder
        if note_shape is not None and note_shape.has_text_frame:
            return note_shape
    except Exception:
        pass

    try:
        notes_frame = notes_slide.notes_text_frame
    except Exception:
        return None
    if notes_frame is None:
        return None

    try:
        for shape in notes_slide.shapes:
            if not getattr(shape, 'has_text_frame', False):
                continue
            try:
                if shape.text_frame._txBody is notes_frame._txBody:
                    return shape
            except Exception:
                continue
    except Exception:
        pass
    return None


def _extract_notes_text(notes_slide) -> str:
    """提取备注文本，并将 PPT 列表语义转换为 Markdown 列表前缀。"""
    notes_frame = notes_slide.notes_text_frame
    if notes_frame is None:
        return ''

    notes_shape = _get_notes_text_shape(notes_slide)
    lines = []
    ordered_counters = {}

    for para in notes_frame.paragraphs:
        has_text = para.text.strip() != ''
        has_math = any(c.tag.split('}')[-1] == 'm' for c in para._p)
        if not has_text and not has_math:
            if lines and lines[-1] != '':
                lines.append('')
            ordered_counters = {}
            continue

        text = ''.join(run.text for run in get_text_runs(para)).strip()
        if not text:
            continue

        bullet_type = get_paragraph_bullet_type(para, notes_shape)
        level = max(int(getattr(para, 'level', 0) or 0), 0)

        if bullet_type == 'numbered':
            ordered_counters[level] = ordered_counters.get(level, 0) + 1
            for deeper_level in [k for k in ordered_counters if k > level]:
                ordered_counters.pop(deeper_level, None)
            line = f'{"  " * level}{ordered_counters[level]}. {text}'
        elif bullet_type == 'bullet':
            ordered_counters = {}
            line = f'{"  " * level}* {text}'
        else:
            ordered_counters = {}
            line = text

        lines.append(line)

    while lines and lines[0] == '':
        lines.pop(0)
    while lines and lines[-1] == '':
        lines.pop()

    return '\n'.join(lines)


def process_picture(config: ConversionConfig, shape, slide_idx) -> Union[ImageElement, None]:
    if config.disable_image:
        return None

    global picture_count

    file_prefix = ''.join(os.path.basename(config.pptx_path).split('.')[:-1])
    pic_name = file_prefix + f'_{picture_count}'
    # python-pptx 的 shape.image.ext 依赖 PIL 识别格式；对于 WMF 等格式可能抛出 UnidentifiedImageError。
    # 这里优先从 image part 的 partname（例如 /ppt/media/image65.wmf）推断扩展名，兜底再用 shape.image.ext。
    img_blob = None
    pic_ext = None
    try:
        pic_ext = shape.image.ext
        img_blob = shape.image.blob
    except Exception:
        try:
            blip = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            rid = None
            if blip is not None:
                rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if rid:
                part = shape.part.related_part(rid)
                partname = str(getattr(part, 'partname', ''))
                _, ext = os.path.splitext(partname)
                pic_ext = ext.lstrip('.').lower() if ext else None
                img_blob = part.blob
        except Exception:
            # 留给下面兜底逻辑
            pic_ext = None
            img_blob = None
    if not pic_ext or img_blob is None:
        # 最后的兜底：保持旧行为（可能仍会失败，但至少错误更明确）
        pic_ext = shape.image.ext
        img_blob = shape.image.blob
    if not os.path.exists(config.image_dir):
        os.makedirs(config.image_dir)

    output_path = config.image_dir / f'{pic_name}.{pic_ext}'
    common_path = os.path.commonpath([config.output_path, config.image_dir])
    img_outputter_path = os.path.relpath(output_path, common_path)
    with open(output_path, 'wb') as f:
        f.write(img_blob)
        picture_count += 1

    # 普通图片
    if pic_ext != 'wmf':
        return ImageElement(path=img_outputter_path, width=config.image_width)

    # WMF 图片：如果 disable_wmf 为 True，则保留原始 WMF 不转换
    if config.disable_wmf:
        logger.info(f'WMF image {output_path} in slide {slide_idx} kept as-is (disable_wmf=True).')
        return ImageElement(path=img_outputter_path, width=config.image_width, original_ext='wmf')

    # WMF 图片：尝试转换为光栅图（PNG/JPG），失败则输出原始格式
    raster_ext = _env_str("PPTX2MD_WMF_RASTER_EXT", "png").lower()
    if raster_ext not in ("png", "jpg", "jpeg"):
        raster_ext = "png"
    dpi = _env_int("PPTX2MD_WMF_DPI", 600)
    jpg_quality = _env_int("PPTX2MD_WMF_JPEG_QUALITY", 92)
    raster_output_path = os.path.splitext(output_path)[0] + f'.{raster_ext}'
    raster_outputter_path = os.path.splitext(img_outputter_path)[0] + f'.{raster_ext}'

    try:
        try:
            Image.open(output_path).save(raster_output_path)
            return ImageElement(path=raster_outputter_path, width=config.image_width)
        except Exception:  # Pillow 转换失败，尝试 wand(ImageMagick) / magick CLI / PowerPoint COM
            # 1) wand（依赖 ImageMagick 且支持 WMF）
            try:
                from wand.image import Image as WandImage  # type: ignore
                try:
                    with WandImage(filename=str(output_path), resolution=dpi) as img:
                        img.format = raster_ext
                        if raster_ext in ("jpg", "jpeg"):
                            img.compression_quality = jpg_quality
                        img.save(filename=raster_output_path)
                    logger.info(f'WMF image {output_path} in slide {slide_idx} converted to {raster_ext} via wand.')
                    return ImageElement(path=raster_outputter_path, width=config.image_width)
                except TypeError:
                    # 某些 wand 版本不支持 resolution 参数
                    with WandImage(filename=str(output_path)) as img:
                        img.format = raster_ext
                        if raster_ext in ("jpg", "jpeg"):
                            img.compression_quality = jpg_quality
                        img.save(filename=raster_output_path)
                    logger.info(f'WMF image {output_path} in slide {slide_idx} converted to {raster_ext} via wand.')
                    return ImageElement(path=raster_outputter_path, width=config.image_width)
            except Exception:
                pass

            # 2) ImageMagick CLI（magick）
            if _convert_wmf_via_magick(str(output_path), raster_output_path, dpi=dpi, jpg_quality=jpg_quality):
                logger.info(f'WMF image {output_path} in slide {slide_idx} converted to {raster_ext} via magick.')
                return ImageElement(path=raster_outputter_path, width=config.image_width)

            # 3) PowerPoint COM：整页导出后按 shape 边界裁剪（高清、最稳）
            if _should_try_wmf_com_fallback():
                try:
                    if _convert_wmf_via_powerpoint_slide_export(
                        config,
                        shape,
                        slide_idx,
                        raster_output_path,
                    ):
                        logger.info(
                            f'WMF image {output_path} in slide {slide_idx} converted to {raster_ext} via PowerPoint COM.'
                        )
                        return ImageElement(path=raster_outputter_path, width=config.image_width)
                except Exception:
                    pass
    except Exception:
        pass

    try:
        logger.warning(
            f'Cannot convert wmf image {output_path} in slide {slide_idx} to {raster_ext}, keeping original wmf.'
        )
        return ImageElement(path=img_outputter_path, width=config.image_width, original_ext='wmf')
    except Exception:
        return ImageElement(path=img_outputter_path, width=config.image_width, original_ext='wmf')


def process_table(config: ConversionConfig, shape, slide_idx) -> Union[TableElement, None]:
    table = [[sum([get_text_runs(p)
                   for p in cell.text_frame.paragraphs], [])
              for cell in row.cells]
             for row in shape.table.rows]
    if len(table) > 0:
        return TableElement(content=table)
    return None


_NS_MC = '{http://schemas.openxmlformats.org/markup-compatibility/2006}'
_NS_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'


def _is_ole_equation_choice(choice) -> bool:
    """判断 mc:Choice 是否为 OLE 公式对象（如 MathType / Equation Editor）。

    经验规则：Choice 内包含 p:oleObj，且 progId 以 "Equation." 开头。
    这样可以避免影响其他 AlternateContent（例如 a14:m 公式）仍按 Choice 展开。
    """
    try:
        ole_obj = choice.find(f'.//{_NS_P}oleObj')
        if ole_obj is None:
            return False
        prog_id = ole_obj.get('progId') or ''
        return prog_id.startswith('Equation.')
    except Exception:
        return False


def _unwrap_alternate_content(slide) -> int:
    """将 slide XML 树中的 mc:AlternateContent 元素展开为其 mc:Choice 内的子元素。

    PowerPoint 在包含公式的 shape 外包裹 mc:AlternateContent，导致 python-pptx 的
    slide.shapes 无法识别它们。此函数在 python-pptx 解析之前修改内存中的 XML 树，
    用 mc:Choice 下的实际 shape 元素（p:sp、p:pic 等）替换 AlternateContent 包裹。

    递归处理 spTree 及其内部的 GroupShape，确保嵌套在组合形状中的
    AlternateContent 也被正确展开。

    仅修改内存 XML，不影响原始 PPTX 文件。

    返回被展开的 AlternateContent 数量。
    """
    sp_tree = slide._element.find(f'.//{_NS_P}cSld/{_NS_P}spTree')
    if sp_tree is None:
        return 0
    # 收集整棵树中所有 AC 元素（包括 GroupShape 内部），
    # reversed() 从深到浅处理，避免修改树时影响尚未处理的元素
    ac_elements = list(sp_tree.iter(f'{_NS_MC}AlternateContent'))
    if not ac_elements:
        return 0
    unwrapped = 0
    for ac in reversed(ac_elements):
        # 元素可能已因父节点被替换而脱离文档，需跳过
        if ac.getparent() is None:
            continue
        choice = ac.find(f'{_NS_MC}Choice')
        if choice is None:
            continue
        # 对 OLE 公式（MathType/Equation Editor）：
        # 预览图通常在 mc:Fallback 的 p:pic 里；若仍按 Choice 展开，会丢失图片引用。
        if _is_ole_equation_choice(choice):
            fallback = ac.find(f'{_NS_MC}Fallback')
            if fallback is not None:
                fallback_pics = list(fallback.iter(f'{_NS_P}pic'))
                if fallback_pics:
                    # OLE 方程的 AlternateContent 经常嵌在 p:graphicFrame/a:graphicData 内部。
                    # 仅替换 ac 本身会把 p:pic 留在 graphicData 里，python-pptx 仍会把它当作 OLE 对象。
                    # 因此尽量用 fallback 的 p:pic 替换其祖先 p:graphicFrame，变成真正的图片 shape。
                    gf = ac
                    while gf is not None and gf.tag != f'{_NS_P}graphicFrame':
                        gf = gf.getparent()
                    if gf is not None and gf.getparent() is not None:
                        gf_parent = gf.getparent()
                        gf_idx = list(gf_parent).index(gf)
                        gf_parent.remove(gf)
                        for i, pic_elem in enumerate(fallback_pics):
                            gf_parent.insert(gf_idx + i, pic_elem)
                        unwrapped += 1
                        continue
                    children = fallback_pics
                else:
                    children = list(choice)
            else:
                children = list(choice)
        else:
            children = list(choice)
        if not children:
            continue
        parent = ac.getparent()
        if parent is None:
            continue
        idx = list(parent).index(ac)
        parent.remove(ac)
        for i, child_elem in enumerate(children):
            parent.insert(idx + i, child_elem)
        unwrapped += 1
    if unwrapped:
        logger.debug(f'Unwrapped {unwrapped} AlternateContent elements in slide')
    return unwrapped


def ungroup_shapes(shapes) -> List[SlideElement]:
    res = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                res.extend(ungroup_shapes(shape.shapes))
            else:
                res.append(shape)
        except Exception as e:
            logger.warning(f'failed to load shape {shape}, skipped. error: {e}')
    return res


def process_shapes(config: ConversionConfig, current_shapes, slide_id: int) -> List[SlideElement]:
    results = []
    for shape in current_shapes:
        if is_title(shape):
            results.append(process_title(config, shape, slide_id))
        elif is_text_block(config, shape):
            results.extend(process_text_blocks(config, shape, slide_id))
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                pic = process_picture(config, shape, slide_id)
                if pic:
                    results.append(pic)
            except AttributeError as e:
                logger.warning(f'Failed to process picture, skipped: {e}')
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = process_table(config, shape, slide_id)
            if table:
                results.append(table)
        else:
            try:
                ph = shape.placeholder_format
                if ph.type == PP_PLACEHOLDER.OBJECT and hasattr(shape, "image") and getattr(shape, "image"):
                    pic = process_picture(config, shape, slide_id)
                    if pic:
                        results.append(pic)
            except:
                pass

    return results


def parse(config: ConversionConfig, prs: Presentation, progress_callback=None, cancel_event=None,
           disable_tqdm=False) -> ParsedPresentation:
    """将 PowerPoint 演示文稿解析为结构化数据。

    参数:
        config: 转换配置。
        prs: PowerPoint 演示文稿对象。
        progress_callback: 可选的进度更新回调，签名: (current, total, slide_name)。
        cancel_event: 可选的 threading.Event，用于支持取消操作。
        disable_tqdm: 禁用 tqdm 进度条（适用于 GUI）。

    返回:
        包含解析后幻灯片数据的 ParsedPresentation。
    """
    result = ParsedPresentation(slides=[])
    slides = list(prs.slides)
    total_slides = len(slides)
    iterator = slides if disable_tqdm else tqdm(slides, desc='Converting slides')
    multi_column_slide_getter = None
    if config.try_multi_column:
        from pptx2md.multi_column import (
            get_multi_column_slide_if_present as multi_column_slide_getter,
        )

    for idx, slide in enumerate(iterator):
        if cancel_event and cancel_event.is_set():
            logger.warning('Conversion cancelled by user.')
            break

        if progress_callback:
            progress_callback(idx + 1, total_slides, f'Slide {idx + 1}')
        if config.page is not None and idx + 1 != config.page:
            continue
        _unwrap_alternate_content(slide)
        shapes = []
        try:
            shapes = sorted(ungroup_shapes(slide.shapes), key=attrgetter('top', 'left'))
        except:
            logger.warning('Bad shapes encountered in this slide. Please check or remove them and try again.')
            logger.warning('shapes:')
            try:
                for sp in slide.shapes:
                    logger.warning(sp.shape_type)
                    logger.warning(sp.top, sp.left, sp.width, sp.height)
            except:
                logger.warning('failed to print all bad shapes.')

        if not config.try_multi_column:
            result_slide = GeneralSlide(elements=process_shapes(config, shapes, idx + 1))
        else:
            multi_column_slide = multi_column_slide_getter(
                prs, slide, partial(process_shapes, config=config, slide_id=idx + 1))
            if multi_column_slide:
                result_slide = multi_column_slide
            else:
                result_slide = GeneralSlide(elements=process_shapes(config, shapes, idx + 1))

        if not config.disable_notes and slide.has_notes_slide:
            notes_text = _extract_notes_text(slide.notes_slide)
            if notes_text:
                result_slide.notes.append(notes_text)

        result.slides.append(result_slide)

    return result
