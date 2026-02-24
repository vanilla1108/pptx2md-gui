import re

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from pptx2md.entry import convert
from pptx2md.types import ConversionConfig


def _clear_bullet_nodes(ppr):
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        node = ppr.find(qn(tag))
        if node is not None:
            ppr.remove(node)


def _set_numbered_para(para, start_at=None):
    ppr = para._p.get_or_add_pPr()
    _clear_bullet_nodes(ppr)
    auto_num = OxmlElement("a:buAutoNum")
    auto_num.set("type", "arabicPeriod")
    if start_at is not None:
        auto_num.set("startAt", str(start_at))
    ppr.append(auto_num)


def _build_and_convert(tmp_path, fill_body):
    pptx_path = tmp_path / "ordered_start_case.pptx"
    output_path = tmp_path / "result.md"
    image_dir = tmp_path / "img"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Ordered StartAt Test"

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    fill_body(body)
    prs.save(pptx_path)

    config = ConversionConfig(
        pptx_path=pptx_path,
        output_path=output_path,
        image_dir=image_dir,
        disable_image=True,
        disable_color=True,
        disable_escaping=True,
        disable_notes=True,
        enable_slides=False,
        min_block_size=3,
    )
    convert(config, disable_tqdm=True)
    return output_path.read_text(encoding="utf-8")


def test_body_ordered_list_uses_start_at(tmp_path):
    def _fill(body):
        p0 = body.paragraphs[0]
        p0.text = "这是第2题"
        _set_numbered_para(p0, start_at=2)

    output = _build_and_convert(tmp_path, _fill)
    assert re.search(r"\n2\. 这是第2题\n", output), output
    assert not re.search(r"\n1\. 这是第2题\n", output), output


def test_body_ordered_list_continues_after_start_at_seed(tmp_path):
    def _fill(body):
        p0 = body.paragraphs[0]
        p0.text = "这是第2题"
        _set_numbered_para(p0, start_at=2)

        p1 = body.add_paragraph()
        p1.text = "这是第3题"
        _set_numbered_para(p1)

    output = _build_and_convert(tmp_path, _fill)
    assert re.search(r"\n2\. 这是第2题\n3\. 这是第3题\n", output), output


def test_body_ordered_list_repeated_same_start_at_still_increments(tmp_path):
    def _fill(body):
        items = ["目录A", "目录B", "目录C", "目录D"]
        for idx, text in enumerate(items):
            para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            para.text = text
            _set_numbered_para(para, start_at=3)

    output = _build_and_convert(tmp_path, _fill)
    assert re.search(r"\n3\. 目录A\n4\. 目录B\n5\. 目录C\n6\. 目录D\n", output), output
