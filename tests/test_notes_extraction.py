import re

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement

from pptx2md.entry import convert
from pptx2md.types import ConversionConfig


def _clear_bullet_nodes(ppr):
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        node = ppr.find(qn(tag))
        if node is not None:
            ppr.remove(node)


def _set_bullet_para(para):
    ppr = para._p.get_or_add_pPr()
    _clear_bullet_nodes(ppr)
    bullet = OxmlElement("a:buChar")
    bullet.set("char", "•")
    ppr.append(bullet)


def _set_numbered_para(para, start_at=1):
    ppr = para._p.get_or_add_pPr()
    _clear_bullet_nodes(ppr)
    auto_num = OxmlElement("a:buAutoNum")
    auto_num.set("type", "arabicPeriod")
    auto_num.set("startAt", str(start_at))
    ppr.append(auto_num)


def _build_and_convert(tmp_path, fill_notes):
    pptx_path = tmp_path / "notes_case.pptx"
    output_path = tmp_path / "result.md"
    image_dir = tmp_path / "img"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Notes Test"

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    fill_notes(notes_frame)
    prs.save(pptx_path)

    config = ConversionConfig(
        pptx_path=pptx_path,
        output_path=output_path,
        image_dir=image_dir,
        disable_image=True,
        disable_color=True,
        disable_escaping=True,
        disable_notes=False,
        enable_slides=False,
        min_block_size=3,
    )
    convert(config, disable_tqdm=True)
    return output_path.read_text(encoding="utf-8")


def test_notes_bullet_list_keeps_markdown_prefix(tmp_path):
    def _fill(notes_frame):
        p0 = notes_frame.paragraphs[0]
        p0.text = "分析问题定位任务"
        _set_bullet_para(p0)

        p1 = notes_frame.add_paragraph()
        p1.text = "数据清理"
        _set_bullet_para(p1)

    output = _build_and_convert(tmp_path, _fill)
    assert re.search(r"\n\* 分析问题定位任务\n\* 数据清理\n", output), output


def test_notes_numbered_list_keeps_markdown_prefix(tmp_path):
    def _fill(notes_frame):
        p0 = notes_frame.paragraphs[0]
        p0.text = "定位任务"
        _set_numbered_para(p0, start_at=1)

        p1 = notes_frame.add_paragraph()
        p1.text = "模型训练"
        _set_numbered_para(p1, start_at=2)

    output = _build_and_convert(tmp_path, _fill)
    assert re.search(r"\n1\. 定位任务\n2\. 模型训练\n", output), output
